from pathlib import Path
from typing import List

from bs4 import BeautifulSoup
import markdown as md
from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation


def _read_pdf(path: str) -> str:
    reader = PdfReader(path)
    texts: List[str] = []
    for page in reader.pages:
        texts.append(page.extract_text() or "")
    return "\n".join(texts)


def _read_docx(path: str) -> str:
    d = DocxDocument(path)
    return "\n".join(p.text for p in d.paragraphs)


def _read_pptx(path: str) -> str:
    prs = Presentation(path)
    texts: List[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
    return "\n".join(texts)


def _read_md(path: str) -> str:
    raw = Path(path).read_text(encoding="utf-8", errors="ignore")
    html = md.markdown(raw)
    soup = BeautifulSoup(html, "html.parser")
    return soup.get_text("\n")


def extract_text_from_path(path: str) -> str:
    suffix = Path(path).suffix.lower()
    if suffix == ".pdf":
        return _read_pdf(path)
    if suffix == ".docx":
        return _read_docx(path)
    if suffix == ".pptx":
        return _read_pptx(path)
    if suffix == ".md":
        return _read_md(path)
    # fallback: plaintext
    return Path(path).read_text(encoding="utf-8", errors="ignore")


